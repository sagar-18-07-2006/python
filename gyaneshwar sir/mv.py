from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# --- Slide 1: Title Slide ---
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "LMV–HMV Detector"
subtitle.text = "Project Introduction\nAutomatic Classification of Light & Heavy Motor Vehicles"

# --- Slide 2: Introduction ---
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction"
content.text = (
    "• Traffic management and road safety are critical challenges today.\n"
    "• Classifying vehicles into Light Motor Vehicles (LMVs) and Heavy Motor Vehicles (HMVs) "
    "is essential for toll collection, traffic regulation, and smart transport systems.\n"
    "• Traditional methods are manual and inefficient, leading to delays and errors.\n"
    "• Among multiple topics, we chose LMV–HMV Detection due to its strong real-world impact.\n"
    "• Objective: Build an automated system that detects and classifies vehicles as LMV or HMV in real-time."
)

# --- Slide 3: Objectives (Optional in intro) ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Objectives"
content.text = (
    "• Develop an automated system for vehicle detection.\n"
    "• Accurately classify vehicles as LMV or HMV.\n"
    "• Enhance efficiency in toll booths, traffic monitoring, and road safety.\n"
    "• Reduce human error and support smart city initiatives."
)

# Save the file
# Correct save path
pptx_file = r"c:\Users\ASUS\OneDrive\Documents\coding\python\gyaneshwar sir\LMV_HMV_Introduction.pptx"
prs.save(pptx_file)

